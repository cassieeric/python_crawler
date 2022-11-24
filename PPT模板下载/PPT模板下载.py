# -*- coding: utf-8 -*-
# @Time    : 2022/5/18 0018 17:47
# @Author  : Thresh
# @File    : demo.py
#
#                             _ooOoo_
#                            o8888888o
#                            88" . "88
#                            (| -_- |)
#                            O\  =  /O
#                         ____/`---'\____
#                       .'  \\|     |//  `.
#                      /  \\|||  :  |||//  \
#                     /  _||||| -:- |||||-  \
#                     |   | \\\  -  /// |   |
#                     | \_|  ''\---/''  |   |
#                     \  .-\__  `-`  ___/-. /
#                   ___`. .'  /--.--\  `. . __
#                ."" '<  `.___\_<|>_/___.'  >'"".
#               | | :  `- \`.;`\ _ /`;.`/ - ` : | |
#               \  \ `-.   \_ __\ /__ _/   .-` /  /
#          ======`-.____`-.___\_____/___.-`____.-'======
#                             `=---='
#         ^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^
#                   ☀佛祖保佑☀   卍   ☀永无BUG☀
import requests, os, zipfile, time
from lxml import etree
from pptx import Presentation
# from unrar import rarfile
from fake_useragent import UserAgent


def get_fenlei_url(url):
    """
    获取每个分类的url
    :param url:
    :return:
    """
    res = requests.get(url, headers={'User-Agent': UserAgent().random})
    res.encoding = 'gb2312'
    et = etree.HTML(res.text)
    uls = et.xpath('//div[@class="col_nav i_nav clearfix"]/ul')  # position()>1 跳过第一个li标签（重点）
    fenlei_urls = {}
    for ul in uls[:-2]:
        for li in ul[1:]:
            title = li.xpath('./a/@title')[0]
            fenlei_url = 'https://www.1ppt.com' + li.xpath('./a/@href')[0]
            fenlei_urls[title] = fenlei_url
    return fenlei_urls


def get_fenlei_allurl(url):
    """
    获取每个分类的所有页面的url
    :param url:
    :return:
    """
    page_num_urls = []
    fenlei_name = url.split('/')[-2]
    page_content = requests.get(url, headers={'User-Agent': UserAgent().random}).text
    et1 = etree.HTML(page_content)
    page_num = len(et1.xpath('//ul[@class="pages"]/li')[1:-2])
    if page_num <= 1:
        page_num_urls.append(url)
    else:
        for i in range(1, 3):  # 这里控制爬取的页码
            all_fenlei_url = url + 'ppt_' + fenlei_name + f'_{i}.html'
            page_num_urls.append(all_fenlei_url)
    return page_num_urls


def get_title_Download_url(url):
    """
    通过详情页解析到标题，并且解析到下载地址请求到真的下载地址
    :param url:
    :return:
    """
    # url = 'https://www.1ppt.com/article/92521.html'
    content = requests.get(url, headers={'User-Agent': UserAgent().random})
    content.encoding = 'gb2312'
    et = etree.HTML(content.text)
    title = et.xpath("//h1/text()")[0]
    down_url = 'https://www.1ppt.com' + et.xpath('//ul[@class="downurllist"]/li/a/@href')[0]
    res = requests.get(down_url, headers={'User-Agent': UserAgent().random}).text
    et_down = etree.HTML(res)
    Download_url = et_down.xpath('//li[@class="c1"]/a/@href')[0]
    old_name = Download_url.split('/')[-1]
    # print(title, Download_url)
    return title, Download_url, old_name


def get_detail_url(folder, url):
    """
    获取每个PPT模板的详情页url和每个ppt模板里的下载地址和标题
    :param folder:
    :param url:
    :return:
    """
    if not os.path.exists('data/' + folder):
        os.mkdir('data/' + folder)
    response = requests.get(url, headers={'User-Agent': UserAgent().random})
    response.encoding = 'utf-8'
    et = etree.HTML(response.text)
    lis = et.xpath('//ul[@class="tplist"]/li')
    detail_urls = []
    for li in lis:
        detail_url = 'https://www.1ppt.com' + li.xpath('./a/@href')[0]
        detail_urls.append(detail_url)
    return detail_urls


def Download_ppt_file(folder_path, old_name, Download_url, title):
    """
    开始下载pptx文件！并根据大小提取压缩文件中指定的pptx文件后删除压缩文件
    :param folder_path:
    :param filename:
    :param url:
    :return:
    """
    path = folder_path + old_name
    with open(path, 'wb') as fp:
        fp.write(requests.get(Download_url, headers={'User-Agent': UserAgent().random}).content)
    print("开始处理压缩文件！")
    try:
        # 判断是什么类型的压缩文件
        if old_name.endswith('.zip'):
            with zipfile.ZipFile(path) as zf:
                for file in zf.namelist():
                    if zf.getinfo(file).file_size > 1024 * 10:  # 1024*1024 = 1MB
                        zf.extract(file, path=folder_path)  # 提取单个文件是extract， 提取所有文件是extractall
                        os.rename(folder_path + file,
                                  folder_path + title + str(time.time()).replace('.', '-') + '.pptx')  # 重命名提取出来的文件
        elif old_name.endswith('.rar'):
            pass
            # with rarfile.RarFile(path) as zf:
            #     for file in zf.namelist():
            #         if zf.getinfo(file).file_size > 1024 * 10:  # 1024*1024 = 1MB
            #             zf.extract(file, path=folder_path)  # 提取单个文件是extract， 提取所有文件是extractall
            #             os.rename(folder_path + file,
            #                       folder_path + title + str(time.time()).replace('.', '-') + '.pptx')  # 重命名提取出来的文件
        os.remove(path)  # 删除压缩包
        print("压缩文件已处理成功！")
    except(FileExistsError) as e:
        print(f"压缩文件处理失败！{e}")
        pass


def del_ppt_max_page(folder_path):
    """
    删除pptx文件最后一页！
    :param folder_path:
    :return:
    """
    names = os.listdir(folder_path)
    for name in names:
        file = folder_path + name
        try:
            prs = Presentation(file)
            # 查看一共几页
            slides = prs.slides
            number_pages = len(slides)
            print('file={},总页数={}'.format(file, number_pages))
            # 删除最后一页
            rId = prs.slides._sldIdLst[-1].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[-1]
            # 保存新的ppt
            prs.save(file)
        except:
            print(f"{name}文件不可处理，请手动处理！")
            continue


def main():
    """
    主函数逻辑！
    :return:
    """
    url = 'https://www.1ppt.com/'
    fenlei_urls = get_fenlei_url(url)
    # print(fenlei_urls)
    for folder, fenlei_url in fenlei_urls.items():
        folder_path = 'data/' + folder + '/'
        print(f'当前将要处理的是{folder}类型的模板！')
        page_num_urls = get_fenlei_allurl(fenlei_url)
        print(f"{folder}类型模板的所有页码的url已经全部获取！{folder}类型模板一共{len(page_num_urls)}页！")
        for page_num_url in page_num_urls:
            detaile_urls = get_detail_url(folder, page_num_url)
            for detaile_url in detaile_urls:
                title, Download_url, old_name = get_title_Download_url(detaile_url)
                print(f'当前要下载的ppt文件的网址是：{Download_url}')
                Download_ppt_file(folder_path, old_name, Download_url, title)
                print(f"{title} ppt文件下载成功！")
            print("开始处理pptx文件！")
            del_ppt_max_page(folder_path)


if __name__ == '__main__':
    main()
